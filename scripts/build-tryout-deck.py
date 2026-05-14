"""Generate TRYOUT 민간실증 발표용 PPTX (한 슬라이드당 한 주제)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title: str, subtitle: str = "") -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(1.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        p.alignment = PP_ALIGN.CENTER
        if subtitle:
            b2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12), Inches(2))
            p2 = b2.text_frame.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(0x44, 0x44, 0x55)
            p2.alignment = PP_ALIGN.CENTER

    def add_content_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12), Inches(0.85))
        tfp = tb.text_frame.paragraphs[0]
        tfp.text = title
        tfp.font.size = Pt(28)
        tfp.font.bold = True
        tfp.font.color.rgb = RGBColor(0x16, 0x24, 0x5C)
        body = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.8), Inches(5.8))
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.space_after = Pt(10)

    add_title_slide(
        "야구장(인천) 관람 동선·혼잡 안내 모바일 웹 실증",
        "실시간 비대면 안내 「지금여기」 | 이앤엑스(ENX)\n"
        "협력파트너: SSG랜더스 (Smart-X Sports)\n"
        "2026 인천스타트업파크 TRYOUT 민간 실증사업",
    )

    add_content_slide(
        "기업 및 서비스 개요",
        [
            "수행기관: 이앤엑스(ENX) — 위치·이동 맥락 기반 모바일 웹 서비스",
            "서비스명: 지금여기 — 앱 설치 없이 QR·URL로 접속",
            "강점: 지도 기반 안내, 이벤트·편의 정보 통합, 이용 로그 기반 운영 개선",
        ],
    )

    add_content_slide(
        "현장 과제 (As-Is)",
        [
            "입장·편의·이벤트 정보가 앱·SNS·현장 안내에 분산 → 통합 확인 어려움",
            "혼잡 시간대 실시간 안내 인프라 부족 → 특정 구역 집중·체류·만족도 이슈",
            "관람객 이용·반응 데이터 수집 체계 미흡 → 개선 근거 확보 어려움",
        ],
    )

    add_content_slide(
        "실증 방향 (To-Be)",
        [
            "모바일 웹 한 화면에서 동선·편의·이벤트 정보 즉시 제공",
            "시간대·구역별 안내로 동선 분산·혼잡 완화 유도",
            "QR 접속 중심으로 설치 장벽 제거, 로그·설문으로 성과·만족도 측정",
        ],
    )

    add_content_slide(
        "실증 개요",
        [
            "실증 장소: 인천 SSG랜더스필드 내·외 관람객 동선 구역 (세부는 파트너 협의 확정)",
            "실증 기간: 협약일 ~ 2026.11.30",
            "사업비: 총 33백만원 (지원 30 / 민간 3)",
            "참여인력: 1명 + 필요 시 외부 협력",
        ],
    )

    add_content_slide(
        "실증 내용 · 핵심 기능",
        [
            "지도·안내 카드: 매점·화장실·출구·이벤트 등 핀·카드 안내",
            "시간대별 맥락 안내: 입장 / 경기 중 / 귀가 동선 등 메시지 전환",
            "이벤트·프로모션 카드 노출, 세션·클릭 등 이용 데이터 수집",
            "만족도 QR 설문 상시 노출",
        ],
    )

    add_content_slide(
        "성과 목표 (정량)",
        [
            "① 누적 접속 세션: 경기·이벤트일 기준 3,000세션 이상 (GA 또는 서버 로그)",
            "② 주요 기능 클릭 전환율: 전체 세션 대비 40% 이상",
            "③ 사용자 만족도 긍정 응답: 응답자 대비 70% 이상 (표본 50명 이상)",
        ],
    )

    add_content_slide(
        "추진 일정 (요약)",
        [
            "5월: 협약·킥오프, 구역·KPI·운영정책 확정",
            "6~7월: 시스템·QR·파일럿 운영 개시",
            "8월: 중간점검, 데이터 분석·피드백 반영",
            "9~10월: 개선·만족도 집중, 안정 운영",
            "11월: 최종 평가·성과보고·정산",
        ],
    )

    add_content_slide(
        "기대 효과",
        [
            "구장·행사 현장형 실시간 안내 운영 모델·템플릿 확보",
            "파트너 연계 데이터 기반 서비스 고도화·후속 과제·제휴 근거",
            "타 구장·축제·지역 행사로 확장 가능한 사업화 방향 검증",
        ],
    )

    out = Path(__file__).resolve().parents[1] / "docs" / "이앤엑스_TRYOUT민간실증_발표자료_2026.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
